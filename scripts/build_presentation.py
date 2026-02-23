"""
Build CS584 Capstone Presentation
Schema-Evolving Knowledge Graphs for Insurance
15-minute presentation deck
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ---------------------------------------------------------------------------
# Color palette — Emory brand-ish blues + gold
# ---------------------------------------------------------------------------
EMORY_BLUE   = RGBColor(0x00, 0x2B, 0x7F)   # #002B7F
EMORY_GOLD   = RGBColor(0xF2, 0xA9, 0x00)   # #F2A900
DARK_GRAY    = RGBColor(0x2D, 0x2D, 0x2D)
MID_GRAY     = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY   = RGBColor(0xF4, 0xF4, 0xF4)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREEN        = RGBColor(0x1B, 0x7F, 0x3B)
RED_DARK     = RGBColor(0xC0, 0x39, 0x2B)
ORANGE       = RGBColor(0xE6, 0x7E, 0x22)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=DARK_GRAY,
                 align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_slide(slide, title, bullets, title_size=28, bullet_size=18,
                     title_color=EMORY_BLUE, sub_bullets=None):
    """Add a header bar + bullet list."""
    # Header bar
    add_rect(slide, 0, 0, 13.33, 1.2, EMORY_BLUE)
    add_text_box(slide, title, 0.3, 0.15, 12.5, 1.0,
                 font_size=title_size, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Gold accent line
    add_rect(slide, 0, 1.2, 13.33, 0.06, EMORY_GOLD)

    # Bullets
    y = 1.5
    for i, bullet in enumerate(bullets):
        # Main bullet
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(y), Inches(12.5), Inches(0.5))
        tb.word_wrap = True
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"  •  {bullet}"
        run.font.size = Pt(bullet_size)
        run.font.color.rgb = DARK_GRAY

        # Sub-bullets if provided
        if sub_bullets and i < len(sub_bullets) and sub_bullets[i]:
            y += 0.42
            for sb in sub_bullets[i]:
                stb = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(12.0), Inches(0.38))
                stb.word_wrap = True
                stf = stb.text_frame
                stf.word_wrap = True
                sp = stf.paragraphs[0]
                srun = sp.add_run()
                srun.text = f"     ‒  {sb}"
                srun.font.size = Pt(bullet_size - 3)
                srun.font.color.rgb = MID_GRAY
                y += 0.36
        else:
            y += 0.44


def code_box(slide, text, left, top, width, height, font_size=11):
    """Monospace code-looking box."""
    add_rect(slide, left, top, width, height, RGBColor(0x1E, 0x1E, 0x1E))
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.1),
        Inches(width - 0.2), Inches(height - 0.15)
    )
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
    run.font.name = "Courier New"
    return txBox


def label_box(slide, text, left, top, width, height, bg=EMORY_BLUE, fg=WHITE, font_size=13):
    add_rect(slide, left, top, width, height, bg)
    add_text_box(slide, text, left + 0.05, top + 0.02, width - 0.1, height - 0.05,
                 font_size=font_size, bold=True, color=fg, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Build slides
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]  # completely blank


# ===== SLIDE 1: TITLE =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, EMORY_BLUE)
add_rect(s, 0, 5.5, 13.33, 2.0, RGBColor(0x00, 0x1A, 0x57))
add_rect(s, 0.4, 4.4, 1.5, 0.12, EMORY_GOLD)

add_text_box(s, "Schema-Evolving Knowledge Graphs\nfor Insurance",
             0.5, 1.0, 12.0, 2.2, font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(s, "Automated Ontology Induction with Cross-Domain Transfer Learning",
             0.5, 3.1, 12.0, 0.8, font_size=22, bold=False, color=RGBColor(0xAA, 0xC4, 0xFF),
             align=PP_ALIGN.LEFT, italic=True)
add_text_box(s, "Xiaofei Wang  ·  Zechary Chou",
             0.5, 4.6, 8.0, 0.5, font_size=18, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(s, "CS 584 AI Capstone  ·  Emory University  ·  Spring 2026",
             0.5, 5.1, 10.0, 0.5, font_size=15, color=RGBColor(0xAA, 0xC4, 0xFF), align=PP_ALIGN.LEFT)
add_text_box(s, "Progress Report — Weeks 1–3",
             0.5, 6.3, 8.0, 0.5, font_size=14, color=EMORY_GOLD, bold=True, align=PP_ALIGN.LEFT)


# ===== SLIDE 2: THE PROBLEM =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 1.2, EMORY_BLUE)
add_text_box(s, "The Problem", 0.3, 0.15, 12.5, 1.0,
             font_size=30, bold=True, color=WHITE)
add_rect(s, 0, 1.2, 13.33, 0.06, EMORY_GOLD)

# Pain point box
add_rect(s, 0.3, 1.5, 12.6, 1.1, RGBColor(0xFF, 0xF0, 0xE0))
add_text_box(s, "⚠  Insurance companies adding a new Line of Business (LOB) face a\n"
                "    3–6 month ontology engineering bottleneck — per domain.",
             0.5, 1.55, 12.2, 1.0, font_size=19, bold=True, color=RGBColor(0x8B, 0x2C, 0x00))

bullets = [
    "Flood, auto, health, liability — each needs its own ontology schema",
    "Current approach: domain experts hand-craft ontologies (slow, expensive, doesn't scale)",
    "No mechanism to reuse structure learned in one domain when adding the next",
    "State-of-practice baseline (LangChain LLMGraphTransformer): extracts triples directly,\n"
    "     no entity resolution, no ontology, no cross-domain transfer",
]
y = 2.85
for b in bullets:
    tb = s.shapes.add_textbox(Inches(0.4), Inches(y), Inches(12.5), Inches(0.65))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"  •  {b}"
    run.font.size = Pt(17)
    run.font.color.rgb = DARK_GRAY
    y += 0.55

add_text_box(s, "Our goal: automate this pipeline — induce ontologies from raw data, extend them across LOBs.",
             0.3, 6.65, 12.5, 0.65, font_size=16, bold=True, color=EMORY_BLUE, italic=True)


# ===== SLIDE 3: PIPELINE ARCHITECTURE =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 1.2, EMORY_BLUE)
add_text_box(s, "Our Approach: 4-Zone Pipeline", 0.3, 0.15, 12.5, 1.0,
             font_size=30, bold=True, color=WHITE)
add_rect(s, 0, 1.2, 13.33, 0.06, EMORY_GOLD)

# Zone boxes
zones = [
    ("Zone 1\nIngestion",       "PDF + CSV\nHybrid chunking\n(section-aware\n+ semantic merge)",  EMORY_BLUE),
    ("Zone 2\nOpen IE",         "Llama 3.1 8B\nFew-shot prompting\n→ (s, r, o) triples\nQuality filtered",  RGBColor(0x15, 0x6E, 0xA0)),
    ("Zone 3\nSchema Evolution","Leiden clustering\nOntology induction\nEntity resolution\nCross-domain transfer", RGBColor(0x1A, 0x7A, 0x4A)),
    ("Zone 4\nNeo4j Storage",   "3-layer graph:\nOntology classes\nEntity types\nEntity instances",  RGBColor(0x6B, 0x35, 0x8A)),
]
box_w = 2.8
x_start = 0.5
for i, (title, body, color) in enumerate(zones):
    x = x_start + i * 3.1
    add_rect(s, x, 1.5, box_w, 4.2, color)
    add_text_box(s, title, x + 0.1, 1.6, box_w - 0.2, 0.75,
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x + 0.1, 2.38, box_w - 0.2, 0.04, WHITE)
    add_text_box(s, body, x + 0.1, 2.5, box_w - 0.2, 3.0,
                 font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow between boxes
    if i < 3:
        add_text_box(s, "→", x + box_w + 0.05, 3.2, 0.35, 0.5,
                     font_size=24, bold=True, color=EMORY_GOLD, align=PP_ALIGN.CENTER)

# Highlight current work
add_rect(s, 0.5, 5.85, 2.8, 0.28, EMORY_GOLD)
add_text_box(s, "✓ DONE", 0.5, 5.87, 2.8, 0.25, font_size=12, bold=True,
             color=EMORY_BLUE, align=PP_ALIGN.CENTER)
add_rect(s, 3.6, 5.85, 2.8, 0.28, EMORY_GOLD)
add_text_box(s, "NEXT →", 3.6, 5.87, 2.8, 0.25, font_size=12, bold=True,
             color=EMORY_BLUE, align=PP_ALIGN.CENTER)

add_text_box(s, "★  Novel contribution: Zone 3 — bottom-up ontology induction via Leiden community detection + cross-domain transfer",
             0.3, 6.3, 12.6, 0.65, font_size=14, bold=True, color=EMORY_BLUE, italic=True)


# ===== SLIDE 4: DATASET =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 1.2, EMORY_BLUE)
add_text_box(s, "Dataset", 0.3, 0.15, 12.5, 1.0, font_size=30, bold=True, color=WHITE)
add_rect(s, 0, 1.2, 13.33, 0.06, EMORY_GOLD)

# Two-column: Flood | Auto
add_rect(s, 0.3, 1.45, 6.0, 4.5, RGBColor(0xE8, 0xF0, 0xFF))
add_text_box(s, "🌊  NFIP Flood Insurance  (Primary)", 0.45, 1.55, 5.7, 0.5,
             font_size=16, bold=True, color=EMORY_BLUE)
flood_items = [
    ("SFIP PDF (F-123)", "27-page policy doc\n       → Zone 1: 49 section-aware chunks\n       (vs. baseline: 56 fixed 512-token chunks)"),
    ("OpenFEMA Policies", "500 records (sample)\n       Full dataset: 100K+"),
    ("OpenFEMA Claims",   "500 records (sample)\n       Full dataset: 50K+"),
    ("Total (full run)", "~150K records across PDF + CSV"),
]
fy = 2.15
for label, detail in flood_items:
    add_text_box(s, f"  •  {label}", 0.45, fy, 5.7, 0.3, font_size=14, bold=True, color=DARK_GRAY)
    add_text_box(s, f"       {detail}", 0.45, fy + 0.28, 5.7, 0.5, font_size=12, color=MID_GRAY)
    fy += 0.82

add_rect(s, 6.8, 1.45, 6.0, 4.5, RGBColor(0xE8, 0xFF, 0xEE))
add_text_box(s, "🚗  Kaggle Auto Insurance  (Transfer study)", 6.95, 1.55, 5.7, 0.5,
             font_size=16, bold=True, color=GREEN)
auto_items = [
    ("60K policy records", "Cross-domain transfer evaluation\n       (Weeks 8–9)"),
    ("Target CTR > 60%",   "Class Transfer Rate — how much\n       flood ontology reuses in auto"),
    ("Target PTR > 50%",   "Property Transfer Rate — shared\n       properties across domains"),
]
ay = 2.15
for label, detail in auto_items:
    add_text_box(s, f"  •  {label}", 6.95, ay, 5.7, 0.3, font_size=14, bold=True, color=DARK_GRAY)
    add_text_box(s, f"       {detail}", 6.95, ay + 0.28, 5.7, 0.5, font_size=12, color=MID_GRAY)
    ay += 0.82

# Reference ontology
add_rect(s, 0.3, 6.1, 12.5, 0.7, RGBColor(0xFF, 0xF8, 0xE1))
add_text_box(s, "📐  Reference Ontology:  Riskine  (github.com/riskine/ontology)  —  used for Precision / Recall / F1 evaluation",
             0.5, 6.17, 12.0, 0.55, font_size=14, color=RGBColor(0x7B, 0x5B, 0x00), bold=True)


# ===== SLIDE 5: BASELINE =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 1.2, EMORY_BLUE)
add_text_box(s, "Baseline: LangChain LLMGraphTransformer", 0.3, 0.15, 12.5, 1.0,
             font_size=28, bold=True, color=WHITE)
add_rect(s, 0, 1.2, 13.33, 0.06, EMORY_GOLD)

add_text_box(s, "The current state-of-practice — what we compare against", 0.4, 1.4, 12.0, 0.45,
             font_size=16, italic=True, color=MID_GRAY)

# Flow diagram
steps = ["Fixed 512-token\nchunks\n(sliding window)", "LLMGraphTransformer\n(Llama 3.1 8B)\nno constraints",
         "Direct Neo4j\ninsertion\n(no deduplication)", "Raw graph\n(no ontology\nno entity resolution)"]
colors_b = [RGBColor(0x4A, 0x6F, 0xA5), RGBColor(0x4A, 0x6F, 0xA5),
            RGBColor(0x4A, 0x6F, 0xA5), RED_DARK]
bx = 0.4
for i, (step, col) in enumerate(zip(steps, colors_b)):
    add_rect(s, bx, 2.1, 2.7, 1.5, col)
    add_text_box(s, step, bx + 0.1, 2.2, 2.5, 1.3, font_size=14, color=WHITE,
                 align=PP_ALIGN.CENTER, bold=(i == 3))
    if i < 3:
        add_text_box(s, "→", bx + 2.75, 2.65, 0.35, 0.5, font_size=22,
                     bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    bx += 3.1

# Known issues
add_text_box(s, "Known weaknesses (what our novel pipeline fixes):", 0.4, 3.85, 12.0, 0.4,
             font_size=15, bold=True, color=RED_DARK)
issues = [
    "20–30% entity duplication   e.g.  'Policy 5123' and 'policy #5123' → two separate nodes",
    "Type inconsistency   e.g.  'Insurance', 'Insurance_policy', 'Insurance coverage' → same concept, 3 labels",
    "No ontology hierarchy   —   flat nodes, no SUBCLASS_OF edges, no class clustering",
    "No cross-domain transfer   —   starts from scratch for every new LOB",
]
iy = 4.3
for issue in issues:
    add_text_box(s, f"  ✗  {issue}", 0.4, iy, 12.5, 0.42, font_size=14, color=DARK_GRAY)
    iy += 0.44


# ===== SLIDE 6: CONCRETE EXAMPLE — CHUNKING =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 1.2, EMORY_BLUE)
add_text_box(s, "Concrete Example 1: Chunking — The Problem", 0.3, 0.15, 12.5, 1.0,
             font_size=28, bold=True, color=WHITE)
add_rect(s, 0, 1.2, 13.33, 0.06, EMORY_GOLD)

# Left — bad chunk
add_rect(s, 0.2, 1.42, 6.1, 0.38, RED_DARK)
add_text_box(s, "✗  BASELINE  — Fixed 512-token chunk (chunk 10, page 5)", 0.25, 1.46, 6.0, 0.32,
             font_size=13, bold=True, color=WHITE)
code_box(s,
    "...time of the loss:\n"
    "a. By over-the-top or frame ties to ground\n"
    "   anchors; or\n"
    "b. In accordance with the manufacturer's\n"
    "   specifications; or\n"
    "c. In compliance with the community's\n"
    "   floodplain management requirements\n"
    "   unless it has been continuously insured\n"
    "   by the NFIP at the same described\n"
    "   location since September 30, 1982.\n"
    "8. Items of property below the lowest\n"
    "   elevated floor of an elevated post-FIRM...",
    0.2, 1.82, 6.1, 3.5, font_size=10.5)
add_text_box(s, "⚠  Mid-sentence start — no context about WHAT rule this applies to.\n"
                "    LLM cannot extract meaningful triples.",
             0.2, 5.35, 6.1, 0.75, font_size=13, color=RED_DARK, bold=True)

# Right — good chunk
add_rect(s, 6.9, 1.42, 6.1, 0.38, GREEN)
add_text_box(s, "✓  ZONE 1  — Section-aware chunk (IV. PROPERTY NOT INSURED)", 6.95, 1.46, 6.0, 0.32,
             font_size=13, bold=True, color=WHITE)
code_box(s,
    "IV. PROPERTY NOT INSURED\n\n"
    "We do not insure any of the following\n"
    "property:\n\n"
    "1. Personal property not inside the\n"
    "   fully enclosed building.\n"
    "2. A building...located entirely in, on,\n"
    "   or over water or seaward of mean\n"
    "   high tide if constructed after\n"
    "   September 30, 1982.\n"
    "3. Open structures, including a building\n"
    "   used as a boathouse or any structure...",
    6.9, 1.82, 6.1, 3.5, font_size=10.5)
add_text_box(s, "✓  Complete section — LLM knows exactly what 'NOT INSURED' means\n"
                "    and can extract (NFIP Policy, NOT_INSURED, Personal property outside building).",
             6.9, 5.35, 6.1, 0.75, font_size=13, color=GREEN, bold=True)

add_text_box(s, "Section hierarchy metadata:  ['IV. PROPERTY NOT INSURED']  |  605 tokens  |  page 11",
             0.3, 6.25, 12.5, 0.45, font_size=13, color=MID_GRAY, italic=True)


# ===== SLIDE 7: CONCRETE EXAMPLE — EXTRACTION OUTPUT =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 1.2, EMORY_BLUE)
add_text_box(s, "Concrete Example 2: Extraction — Type Inconsistency", 0.3, 0.15, 12.5, 1.0,
             font_size=28, bold=True, color=WHITE)
add_rect(s, 0, 1.2, 13.33, 0.06, EMORY_GOLD)

add_text_box(s, 'Input: "Policy" concept appears across 5 different section chunks  →  LLMGraphTransformer produces 5 different node labels:',
             0.3, 1.38, 12.5, 0.5, font_size=15, color=DARK_GRAY)

# The 5 labels
labels_ex = ["Policy", "Policy action", "Policy feature", "Policy type", "Policy clause"]
label_colors = [EMORY_BLUE, RGBColor(0x2E, 0x86, 0xC1), RGBColor(0x1A, 0x5C, 0x8A),
                RGBColor(0x15, 0x4E, 0x7A), RGBColor(0x0D, 0x3F, 0x6B)]
lx = 0.6
for i, (lbl, col) in enumerate(zip(labels_ex, label_colors)):
    add_rect(s, lx, 2.0, 2.2, 0.65, col)
    add_text_box(s, f":{lbl}", lx + 0.1, 2.08, 2.0, 0.5,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(s, "≠", lx + 2.25, 2.15, 0.3, 0.4, font_size=20, bold=True,
                     color=RED_DARK, align=PP_ALIGN.CENTER)
    lx += 2.5

add_text_box(s, "Same real-world concept.  Different chunk → different label.  Graph cannot answer \"What are the policy rules?\"",
             0.3, 2.8, 12.5, 0.45, font_size=15, color=RED_DARK, bold=True)

# Arrows to solution
add_rect(s, 2.0, 3.45, 9.0, 0.06, EMORY_GOLD)
add_text_box(s, "Our Zone 3 Fix: Leiden Clustering", 3.5, 3.6, 6.0, 0.4,
             font_size=17, bold=True, color=EMORY_BLUE, align=PP_ALIGN.CENTER)

code_box(s,
    "# Step 1: Compute type-similarity graph\n"
    "w(Policy, Policy_action) = 0.5 * sem_sim + 0.3 * cooccur + 0.2 * rel_sim\n"
    "                         = 0.5 * 0.91   + 0.3 * 0.72   + 0.2 * 0.68\n"
    "                         = 0.455 + 0.216 + 0.136  =  0.807  → SAME CLUSTER\n\n"
    "# Step 2: Leiden community detection groups all 5 → one cluster\n"
    "C_policy = { Policy, Policy_action, Policy_feature, Policy_type, Policy_clause }\n\n"
    "# Step 3: LLM names the abstract class\n"
    '> "Given these types, generate an abstract class name: InsurancePolicy"',
    0.3, 4.1, 12.6, 2.55, font_size=11)

add_text_box(s, "Result: one canonical :InsurancePolicy class.  All instances unified.  Graph queries now work.",
             0.3, 6.78, 12.5, 0.45, font_size=14, bold=True, color=GREEN)


# ===== SLIDE 8: EVALUATION SETUP =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 1.2, EMORY_BLUE)
add_text_box(s, "Evaluation Method & Metrics", 0.3, 0.15, 12.5, 1.0,
             font_size=30, bold=True, color=WHITE)
add_rect(s, 0, 1.2, 13.33, 0.06, EMORY_GOLD)

# Left: metrics
add_rect(s, 0.2, 1.42, 5.8, 5.3, LIGHT_GRAY)
add_text_box(s, "Automatic Metrics", 0.35, 1.5, 5.5, 0.4,
             font_size=17, bold=True, color=EMORY_BLUE)
metrics_list = [
    ("Query Accuracy (20 tasks)", "% of policy questions answerable\nvia Cypher queries  |  target: >75%"),
    ("Entity Duplication Rate",   "% excess nodes for same entity\n(baseline: ~0%, full data: 20–30%)"),
    ("Type Inconsistency",        "% root concepts with >1 label\n(baseline: 8–15%,  novel: <2%)"),
    ("Precision / Recall / F1",   "vs. Riskine reference ontology\n(semantic match via LLM judge)\ntargets: P>80%, R>70%, F1>0.75"),
]
my = 2.0
for mname, mdetail in metrics_list:
    add_text_box(s, f"  •  {mname}", 0.35, my, 5.5, 0.3, font_size=14, bold=True, color=DARK_GRAY)
    add_text_box(s, f"       {mdetail}", 0.35, my + 0.28, 5.5, 0.6, font_size=12, color=MID_GRAY)
    my += 0.92

# Right: 20 tasks by category
add_rect(s, 6.5, 1.42, 6.5, 5.3, LIGHT_GRAY)
add_text_box(s, "20 Query Tasks — by Category", 6.65, 1.5, 6.2, 0.4,
             font_size=17, bold=True, color=EMORY_BLUE)
categories = [
    ("Coverage (6 tasks)",      ["What property types are covered?", "Max building coverage?", "Basement contents?"]),
    ("Exclusions (3 tasks)",    ["Earth movement covered?", "Financial losses excluded?", "Types of excluded damage?"]),
    ("Definitions (2 tasks)",   ["Define 'flood'?", "Define 'building'?"]),
    ("Claims (3 tasks)",        ["Policyholder obligations after loss?", "Proof of loss deadline?", "Appeal process?"]),
    ("Policy Terms (6 tasks)",  ["Waiting period?", "ICC coverage?", "RCV determination?", "Coinsurance?"]),
]
cy = 2.0
for cat, examples in categories:
    add_text_box(s, f"  •  {cat}", 6.65, cy, 6.2, 0.3, font_size=13, bold=True, color=DARK_GRAY)
    cy += 0.3
    add_text_box(s, "       e.g. " + examples[0], 6.65, cy, 6.2, 0.28, font_size=11, color=MID_GRAY, italic=True)
    cy += 0.6

add_text_box(s, "Manual evaluation: 2 domain experts rate 50 classes + 30 hierarchies (Cohen's κ > 0.7)",
             0.3, 6.85, 12.5, 0.4, font_size=13, color=MID_GRAY, italic=True)


# ===== SLIDE 9: RESULTS TABLE =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 1.2, EMORY_BLUE)
add_text_box(s, "Results: Baseline vs. Zone 1 Ablation", 0.3, 0.15, 12.5, 1.0,
             font_size=30, bold=True, color=WHITE)
add_rect(s, 0, 1.2, 13.33, 0.06, EMORY_GOLD)

# Table
headers = ["Metric", "Baseline\n(512-token chunks)", "Baseline +\nZone 1 Chunks", "Novel Pipeline\n(target)"]
col_widths = [3.3, 2.8, 2.8, 2.8]
h_colors  = [EMORY_BLUE, EMORY_BLUE, RGBColor(0x15, 0x6E, 0xA0), GREEN]

tx = 0.3
for i, (hdr, cw, hc) in enumerate(zip(headers, col_widths, h_colors)):
    add_rect(s, tx, 1.45, cw - 0.05, 0.75, hc)
    add_text_box(s, hdr, tx + 0.05, 1.5, cw - 0.15, 0.65,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx += cw

rows_data = [
    ("Nodes",              "365",    "242 (↓34%)",   "~150–200"),
    ("Unique labels",      "58",     "44 (↓24%)",    "20–30"),
    ("Type inconsistency", "8.0%",   "15.2% (↑)",    "< 2%"),
    ("Query accuracy",     "35%",    "50% (↑15%)",   "> 75%"),
    ("Entity duplication", "~0%",    "~0%",          "< 5%"),
    ("F1 vs Riskine",      "—",      "—",            "> 0.75"),
]
row_alt = [WHITE, LIGHT_GRAY]
ry = 2.22
for ri, (metric, v1, v2, vtarget) in enumerate(rows_data):
    tx = 0.3
    bg = row_alt[ri % 2]
    for ci, (val, cw) in enumerate(zip([metric, v1, v2, vtarget], col_widths)):
        add_rect(s, tx, ry, cw - 0.05, 0.52, bg)
        fc = EMORY_BLUE if ci == 0 else DARK_GRAY
        bold = ci == 0
        # Highlight improvements
        if ci == 2 and "↑" in val:
            fc = GREEN
            bold = True
        if ci == 2 and "15%" in val:
            fc = GREEN
            bold = True
        if ci == 3:
            fc = GREEN
            bold = True
        add_text_box(s, val, tx + 0.08, ry + 0.08, cw - 0.2, 0.38,
                     font_size=14, bold=bold, color=fc, align=PP_ALIGN.CENTER)
        tx += cw
    ry += 0.54

# Note on type inconsistency
add_rect(s, 0.3, 5.58, 12.6, 0.65, RGBColor(0xFF, 0xF8, 0xE1))
add_text_box(s, "★  Type inconsistency went UP with Zone 1 — richer context → more specific labels.\n"
                "    This is expected and correct: Zone 3 Leiden clustering is designed to merge them.",
             0.45, 5.62, 12.3, 0.58, font_size=13, color=RGBColor(0x7B, 0x5B, 0x00), bold=True)

add_text_box(s, "Both baseline runs: same LLMGraphTransformer extractor — only the chunking strategy differs.",
             0.3, 6.35, 12.5, 0.4, font_size=13, color=MID_GRAY, italic=True)


# ===== SLIDE 10: KEY FINDING =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 1.2, EMORY_BLUE)
add_text_box(s, "Key Finding: Zone 1 Chunking Impact", 0.3, 0.15, 12.5, 1.0,
             font_size=30, bold=True, color=WHITE)
add_rect(s, 0, 1.2, 13.33, 0.06, EMORY_GOLD)

# Big number
add_rect(s, 0.3, 1.45, 5.5, 2.1, RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(s, "+15%", 0.5, 1.55, 5.1, 1.1, font_size=72, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text_box(s, "query accuracy gain\nchunking strategy alone", 0.5, 2.65, 5.1, 0.8,
             font_size=16, color=GREEN, align=PP_ALIGN.CENTER)

add_rect(s, 6.3, 1.45, 6.7, 2.1, RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(s, "35% → 50%", 6.5, 1.55, 6.3, 1.1, font_size=52, bold=True, color=EMORY_BLUE, align=PP_ALIGN.CENTER)
add_text_box(s, "baseline → baseline + Zone 1 chunking", 6.5, 2.65, 6.3, 0.8,
             font_size=14, color=EMORY_BLUE, align=PP_ALIGN.CENTER)

# Task breakdown
add_text_box(s, "What changed between the two runs (same LLM, same extraction, different chunks):",
             0.3, 3.75, 12.5, 0.4, font_size=15, bold=True, color=DARK_GRAY)

gained = [
    ("Task 5",  "Earth movement exclusion", "V. EXCLUSIONS kept intact → LLM saw complete exclusion list"),
    ("Task 12", "Building definition",      "II. DEFINITIONS section complete → full definition in one chunk"),
    ("Task 15", "Coinsurance rule",         "VII.B Other Insurance section complete → co-insurance rule visible"),
    ("Task 19", "RCV determination",        "VII.R Loss Settlement complete → replacement cost value captured"),
]
gy = 4.25
for tid, question, reason in gained:
    add_rect(s, 0.3, gy, 0.55, 0.36, GREEN)
    add_text_box(s, "✓", 0.3, gy + 0.02, 0.55, 0.34, font_size=16, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, f"{tid}: {question}", 0.95, gy, 4.0, 0.2, font_size=13, bold=True, color=GREEN)
    add_text_box(s, reason, 0.95, gy + 0.2, 12.0, 0.2, font_size=11, color=MID_GRAY, italic=True)
    gy += 0.48

add_rect(s, 0.3, gy + 0.05, 0.55, 0.36, RED_DARK)
add_text_box(s, "✗", 0.3, gy + 0.07, 0.55, 0.34, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "Task 10: Proof-of-loss deadline", 0.95, gy + 0.05, 5.0, 0.2, font_size=13, bold=True, color=RED_DARK)
add_text_box(s, "Was a lucky keyword hit ('proof of loss') in original; Zone1 chunk returns 1 row without matching keywords",
             0.95, gy + 0.25, 12.0, 0.2, font_size=11, color=MID_GRAY, italic=True)


# ===== SLIDE 11: WHAT BASELINE REVEALS =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 1.2, EMORY_BLUE)
add_text_box(s, "What the Baseline Reveals (Why Our Novel Pipeline Matters)",
             0.3, 0.15, 12.5, 1.0, font_size=26, bold=True, color=WHITE)
add_rect(s, 0, 1.2, 13.33, 0.06, EMORY_GOLD)

problems = [
    (RED_DARK,   "✗  Query accuracy capped at 50%",
     "Even with perfect section-aware chunks, LLMGraphTransformer misses numerical facts\n"
     "(coverage limits), negation (exclusions), and procedural knowledge (claims process)."),
    (RED_DARK,   "✗  Type proliferation — 44 labels for ~10 real concepts",
     "LLM invents new labels per chunk: Policy / Policy action / Policy feature / Policy type / Policy clause.\n"
     "→ Zone 3 Leiden clustering will merge all 5 into one canonical :InsurancePolicy class."),
    (RED_DARK,   "✗  No structured exclusions / negation",
     "Tasks 4, 5, 18 consistently fail — 'NOT INSURED' text is not captured as a typed relationship.\n"
     "→ Zone 2 Open IE will use explicit prompts to extract EXCLUDED_FROM triples."),
    (GREEN,      "✓  Section-aware chunks give LLM coherent context",
     "+15% accuracy from Zone 1 alone confirms that chunking quality directly impacts extraction quality."),
    (GREEN,      "✓  Baseline numbers are our Table 1 comparison row",
     "35% (original) / 50% (+ Zone1) are the documented starting points we need to beat."),
]
py = 1.42
for color, title, detail in problems:
    add_rect(s, 0.2, py, 12.8, 0.9, RGBColor(0xF8, 0xF8, 0xF8))
    add_rect(s, 0.2, py, 0.18, 0.9, color)
    add_text_box(s, title, 0.5, py + 0.05, 12.3, 0.3, font_size=14, bold=True, color=color)
    add_text_box(s, detail, 0.5, py + 0.36, 12.3, 0.48, font_size=12, color=MID_GRAY)
    py += 1.0


# ===== SLIDE 12: NEXT STEPS =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 1.2, EMORY_BLUE)
add_text_box(s, "Next Steps: Novel Pipeline (Weeks 4–7)", 0.3, 0.15, 12.5, 1.0,
             font_size=30, bold=True, color=WHITE)
add_rect(s, 0, 1.2, 13.33, 0.06, EMORY_GOLD)

nexts = [
    ("Zone 2 — Open IE  (next)", RGBColor(0x15, 0x6E, 0xA0),
     [
         "Few-shot prompt Llama 3.1 8B → {subject, relation, object, confidence}",
         "Quality filters: confidence > 0.7 | reject generic rels: has, is, contains",
         "Entity/Event distinction: Policy (entity) vs. Filed (event → typed relation)",
         "Expected: cleaner triples, fewer spurious nodes, better query accuracy",
     ]),
    ("Zone 3 — Ontology Induction", RGBColor(0x1A, 0x7A, 0x4A),
     [
         "Entity resolution: string similarity + embedding cosine (threshold 0.92)",
         "Leiden community detection on type-similarity graph",
         "LLM names abstract classes, infers SUBCLASS_OF hierarchy",
         "Target: 20–30 canonical classes from ~44 raw labels",
     ]),
    ("Zone 3b — Cross-Domain Extension  (Weeks 8–9)", RGBColor(0x6B, 0x35, 0x8A),
     [
         "Induce auto insurance ontology independently from Kaggle data",
         "Match auto classes to flood ontology (DIRECT_REUSE / SUBCLASS / SIBLING / NEW_BRANCH)",
         "Measure CTR > 60% and PTR > 50% — the core research contribution",
     ]),
]
ny = 1.45
for zone_title, zone_color, zone_bullets in nexts:
    add_rect(s, 0.2, ny, 12.8, 0.35, zone_color)
    add_text_box(s, zone_title, 0.35, ny + 0.05, 12.5, 0.27,
                 font_size=15, bold=True, color=WHITE)
    ny += 0.38
    for b in zone_bullets:
        add_text_box(s, f"    •  {b}", 0.4, ny, 12.5, 0.32, font_size=13, color=DARK_GRAY)
        ny += 0.32
    ny += 0.15


# ===== SLIDE 13: TIMELINE =====
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 1.2, EMORY_BLUE)
add_text_box(s, "Timeline", 0.3, 0.15, 12.5, 1.0, font_size=30, bold=True, color=WHITE)
add_rect(s, 0, 1.2, 13.33, 0.06, EMORY_GOLD)

phases = [
    ("Weeks 1–3",   "Baseline Implementation",              "DONE ✓",
     "LangGraph pipeline | Zone 1 chunking | 20-task evaluation | Ablation study",
     GREEN, True),
    ("Weeks 4–7",   "Novel Implementation (OpenFEMA)",       "IN PROGRESS →",
     "Zone 2 Open IE | Zone 3 Leiden clustering | Entity resolution | Full Neo4j graph",
     EMORY_BLUE, False),
    ("Weeks 8–9",   "Cross-Domain Transfer Study",           "UPCOMING",
     "Auto insurance ingestion | CTR/PTR measurement | Ontology merging",
     RGBColor(0x6B, 0x35, 0x8A), False),
    ("Weeks 10–11", "Polishing & Reporting",                 "UPCOMING",
     "Streamlit demo | Final evaluations | 10–15 page report",
     MID_GRAY, False),
]

ty = 1.45
for week, title, status, detail, color, done in phases:
    add_rect(s, 0.2, ty, 12.8, 1.05, RGBColor(0xF8, 0xF8, 0xF8))
    add_rect(s, 0.2, ty, 0.22, 1.05, color)
    # Week label
    add_rect(s, 0.5, ty + 0.12, 1.5, 0.38, color)
    add_text_box(s, week, 0.52, ty + 0.14, 1.46, 0.35, font_size=12, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
    # Title + status
    add_text_box(s, title, 2.15, ty + 0.08, 7.5, 0.35, font_size=16, bold=True, color=DARK_GRAY)
    status_color = GREEN if done else (EMORY_BLUE if "PROGRESS" in status else MID_GRAY)
    add_text_box(s, status, 10.0, ty + 0.08, 2.9, 0.35, font_size=14, bold=True,
                 color=status_color, align=PP_ALIGN.RIGHT)
    add_text_box(s, detail, 2.15, ty + 0.5, 10.5, 0.45, font_size=12, color=MID_GRAY)
    ty += 1.2

add_text_box(s, "Deliverables: LangGraph pipeline  ·  Neo4j graphs (flood / auto / merged)  ·  Streamlit demo  ·  10–15 page report",
             0.3, 6.42, 12.5, 0.4, font_size=13, color=MID_GRAY, italic=True)


# ===== SAVE =====
out_path = "/Users/sam/Documents/School/Emory/CS584_AI_Capstone/CS584_Capstone_Progress_Presentation.pptx"
prs.save(out_path)
print(f"✓ Saved: {out_path}")
print(f"  Slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    shapes = [sh for sh in slide.shapes if sh.has_text_frame]
    first_text = shapes[0].text_frame.text[:60].replace('\n', ' ') if shapes else ""
    print(f"  Slide {i+1}: {first_text}")
